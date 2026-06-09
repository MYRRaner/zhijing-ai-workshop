from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
PRIMARY = RGBColor(0x6C, 0x5C, 0xE7)
PRIMARY_DARK = RGBColor(0x5A, 0x4B, 0xD1)
SECONDARY = RGBColor(0x00, 0xCE, 0xC9)
ACCENT = RGBColor(0xFD, 0x79, 0xA8)
DARK = RGBColor(0x2D, 0x34, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY = RGBColor(0x63, 0x6E, 0x72)
GRADIENT_START = RGBColor(0x0C, 0x0C, 0x1D)
GRADIENT_END = RGBColor(0x2D, 0x1B, 0x69)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = GRADIENT_START
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = GRADIENT_END
    fill.gradient_stops[1].position = 1.0


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox


def add_icon_card(slide, left, top, width, height, icon_text, title, desc, bg_color, icon_color):
    card = add_shape(slide, left, top, width, height, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    icon_shape = add_shape(slide, left + Inches(0.4), top + Inches(0.3), Inches(0.7), Inches(0.7), icon_color)
    icon_tf = icon_shape.text_frame
    icon_tf.paragraphs[0].text = icon_text
    icon_tf.paragraphs[0].font.size = Pt(22)
    icon_tf.paragraphs[0].font.color.rgb = WHITE
    icon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_tf.paragraphs[0].font.bold = True

    add_text_box(slide, left + Inches(0.3), top + Inches(1.15), width - Inches(0.6), Inches(0.4),
                 title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.55), width - Inches(0.6), height - Inches(1.8),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

# 装饰圆形
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x6C, 0x5C, 0xE7)
circle1.fill.fore_color.brightness = 0.8
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x00, 0xCE, 0xC9)
circle2.fill.fore_color.brightness = 0.8
circle2.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "智境 · AI创意工坊", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "融合计算机视觉 · 大语言模型 · AIGC 的智能创作平台", font_size=24, color=RGBColor(0xA0, 0xA0, 0xD0), alignment=PP_ALIGN.CENTER)

# 技术标签
tags = ["Flask", "Django", "OpenCV", "GLM-4", "CogView"]
tag_left = Inches(3.5)
for tag in tags:
    tag_shape = add_shape(slide, tag_left, Inches(4.0), Inches(1.2), Inches(0.45), PRIMARY)
    tag_tf = tag_shape.text_frame
    tag_tf.paragraphs[0].text = tag
    tag_tf.paragraphs[0].font.size = Pt(14)
    tag_tf.paragraphs[0].font.color.rgb = WHITE
    tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_left += Inches(1.4)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Python Web 期末项目", font_size=18, color=RGBColor(0x80, 0x80, 0xA0), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "2026年6月", font_size=16, color=RGBColor(0x80, 0x80, 0xA0), alignment=PP_ALIGN.CENTER)


# ========== 第2页：项目背景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# 顶部色带
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "项目背景与动机", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)

# 左侧：问题
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "痛点分析", font_size=20, color=PRIMARY, bold=True)
pain_items = [
    "传统创作工具缺乏AI辅助，效率低下",
    "AI能力分散，用户需在多个平台间切换",
    "单一AI能力有限，无法实现跨模态创作",
    "缺乏创作社区，作品难以分享和交流",
]
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(3.5), pain_items, font_size=15, color=DARK)

# 右侧：方案
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "解决方案", font_size=20, color=SECONDARY, bold=True)
sol_items = [
    "一站式AI创作平台，集成多种AI能力",
    "微服务架构，AI服务独立部署易扩展",
    "创意工坊模式，组合CV+LLM+AIGC",
    "作品画廊社区，点赞评论形成闭环",
]
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(3.5), sol_items, font_size=15, color=DARK)

# 底部应用场景
add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.4), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(0.4),
             "落地应用场景", font_size=18, color=PRIMARY, bold=True)
scenarios = "内容创作者的智能工具箱 | 设计师的AI灵感助手 | 教育场景的创意教学平台 | 社交媒体的AI内容生成器"
add_text_box(slide, Inches(1.0), Inches(6.15), Inches(11), Inches(0.6),
             scenarios, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ========== 第3页：技术架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "技术架构", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)

# Django 层
django_box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), RGBColor(0xF0, 0xF0, 0xFF), PRIMARY)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
             "Django 主应用 (Port 8000)", font_size=20, color=PRIMARY, bold=True)

django_items = [
    "用户认证系统 (Django Auth)",
    "  - 注册 / 登录 / 个人中心",
    "创意工坊模块 (creative)",
    "  - 智能识物 / AI写作 / AI绘画 / 创意工坊",
    "作品画廊模块 (gallery)",
    "  - 作品展示 / 点赞 / 评论 / 发布",
    "数据库 ORM (SQLite)",
    "  - User / CreativeProject / Artwork / Comment",
]
add_bullet_list(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(3.8), django_items, font_size=13, color=DARK, spacing=Pt(4))

# 箭头
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.5), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = PRIMARY
arrow.line.fill.background()
add_text_box(slide, Inches(6.3), Inches(2.8), Inches(1.2), Inches(0.5),
             "REST API", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.3), Inches(4.1), Inches(1.2), Inches(0.5),
             "JSON", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Flask 层
flask_box = add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(4.8), RGBColor(0xF0, 0xFF, 0xFF), SECONDARY)
add_text_box(slide, Inches(7.7), Inches(1.7), Inches(4.5), Inches(0.5),
             "Flask AI 微服务 (Port 5001)", font_size=20, color=RGBColor(0x00, 0xA0, 0x9C), bold=True)

flask_items = [
    "CV 服务 (OpenCV)",
    "  - 颜色提取 / 人脸检测 / 场景识别",
    "  - 清晰度评估 / 亮度分析",
    "LLM 服务 (GLM-4-Flash)",
    "  - 多风格写作 / 提示词增强",
    "  - 5种风格：通用/文学/技术/幽默/诗歌",
    "AIGC 服务 (CogView-3-Flash)",
    "  - AI图像生成 / 风格迁移",
    "  - 6种风格：动漫/写实/水彩/油画/赛博朋克/国画",
]
add_bullet_list(slide, Inches(7.9), Inches(2.3), Inches(4.3), Inches(3.8), flask_items, font_size=13, color=DARK, spacing=Pt(4))

# 底部说明
add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.5),
             "架构优势：职责分离 | 独立部署 | 易于扩展 | 优雅降级（无API Key时自动Mock演示）",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ========== 第4页：核心功能 - 智能识物 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RGBColor(0x09, 0x84, 0xE3))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "核心功能一：智能识物", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), RGBColor(0x09, 0x84, 0xE3))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "基于 OpenCV 的计算机视觉分析", font_size=16, color=GRAY)

# 功能卡片
cards = [
    ("🎨", "颜色提取", "K-Means聚类算法\n提取5种主色调\n显示HEX值和占比", RGBColor(0x74, 0xB9, 0xFF)),
    ("👤", "人脸检测", "Haar级联分类器\n检测人脸位置\n返回坐标信息", RGBColor(0x09, 0x84, 0xE3)),
    ("🏔", "场景识别", "亮度/清晰度/边缘\n密度综合判断\n场景类型分类", RGBColor(0x00, 0xCE, 0xC9)),
    ("📐", "清晰度评估", "Laplacian方差\n计算图像清晰度\n模糊/锐利判断", RGBColor(0x6C, 0x5C, 0xE7)),
]

card_left = Inches(0.8)
for icon, title, desc, color in cards:
    add_icon_card(slide, card_left, Inches(2.1), Inches(2.8), Inches(3.0), icon, title, desc, WHITE, color)
    card_left += Inches(3.05)

# 流程
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11), Inches(0.4),
             "处理流程", font_size=18, color=DARK, bold=True)
flow_steps = ["上传图片", "→", "OpenCV解码", "→", "K-Means颜色聚类", "→", "Haar人脸检测", "→", "场景/清晰度分析", "→", "返回JSON结果"]
flow_left = Inches(1.0)
for step in flow_steps:
    if step == "→":
        add_text_box(slide, flow_left, Inches(6.2), Inches(0.4), Inches(0.4),
                     "→", font_size=18, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        flow_left += Inches(0.4)
    else:
        step_shape = add_shape(slide, flow_left, Inches(6.15), Inches(1.5), Inches(0.45), RGBColor(0x09, 0x84, 0xE3))
        step_tf = step_shape.text_frame
        step_tf.paragraphs[0].text = step
        step_tf.paragraphs[0].font.size = Pt(11)
        step_tf.paragraphs[0].font.color.rgb = WHITE
        step_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        flow_left += Inches(1.6)


# ========== 第5页：核心功能 - AI写作 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "核心功能二：AI写作", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "基于智谱AI GLM-4-Flash 大语言模型", font_size=16, color=GRAY)

# 风格卡片
styles = [
    ("通用", "自然流畅的通用写作风格", PRIMARY),
    ("文学", "优美富有文学色彩，注重修辞意境", RGBColor(0xA2, 0x9B, 0xFE)),
    ("技术", "专业严谨的技术语言，注重逻辑", RGBColor(0x00, 0xB8, 0x94)),
    ("幽默", "风趣幽默，加入有趣比喻和段子", RGBColor(0xFD, 0xCB, 0x6E)),
    ("诗歌", "诗歌形式，注意韵律和意境", ACCENT),
]

style_left = Inches(0.8)
for name, desc, color in styles:
    card = add_shape(slide, style_left, Inches(2.1), Inches(2.3), Inches(1.8), WHITE, color)
    add_text_box(slide, style_left + Inches(0.15), Inches(2.2), Inches(2.0), Inches(0.4),
                 name, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, style_left + Inches(0.15), Inches(2.7), Inches(2.0), Inches(1.0),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    style_left += Inches(2.45)

# 提示词增强
add_shape(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.6), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(0.4),
             "特色功能：提示词增强（Prompt Engineering）", font_size=18, color=PRIMARY, bold=True)

prompt_items = [
    "创意工坊模式下，LLM可将中文描述自动转化为英文AI绘画提示词",
    "自动添加风格、光影、细节等要素，提升AIGC生成质量",
    "示例：\"一只猫在月光下\" → \"a cat under moonlight, highly detailed, soft lighting, 8k, masterpiece\"",
    "支持多轮对话历史传递，实现上下文连贯的创意写作",
]
add_bullet_list(slide, Inches(1.2), Inches(4.9), Inches(11), Inches(2.0), prompt_items, font_size=14, color=DARK, spacing=Pt(6))


# ========== 第6页：核心功能 - AI绘画 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "核心功能三：AI绘画", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), ACCENT)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "基于智谱AI CogView-3-Flash 图像生成模型", font_size=16, color=GRAY)

# 风格展示
paint_styles = [
    ("动漫", "anime style\nvibrant colors", RGBColor(0xFF, 0x6B, 0x9D)),
    ("写实", "photorealistic\nhigh detail", RGBColor(0x2D, 0x34, 0x36)),
    ("水彩", "watercolor\nsoft edges", RGBColor(0xA8, 0xE6, 0xCF)),
    ("油画", "oil painting\nrich textures", RGBColor(0xE1, 0x70, 0x55)),
    ("赛博朋克", "cyberpunk\nneon lights", RGBColor(0x35, 0x00, 0xD3)),
    ("国画", "Chinese ink\nelegant", RGBColor(0x8B, 0x45, 0x13)),
]

style_left = Inches(0.8)
for name, desc, color in paint_styles:
    card = add_shape(slide, style_left, Inches(2.1), Inches(1.9), Inches(2.2), WHITE, color)
    color_box = add_shape(slide, style_left + Inches(0.3), Inches(2.3), Inches(1.3), Inches(0.8), color)
    add_text_box(slide, style_left + Inches(0.15), Inches(3.2), Inches(1.6), Inches(0.4),
                 name, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, style_left + Inches(0.15), Inches(3.6), Inches(1.6), Inches(0.6),
                 desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)
    style_left += Inches(2.05)

# 生成流程
add_shape(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(0.4),
             "AI绘画生成流程", font_size=18, color=ACCENT, bold=True)

gen_items = [
    "1. 用户输入画面描述（中文/英文均可）",
    "2. 系统自动添加风格后缀（如 \"anime style, vibrant colors, detailed illustration\"）",
    "3. 调用 CogView-3-Flash API 生成图片",
    "4. 返回生成结果，支持下载保存",
    "5. 优雅降级：无API Key时使用Pillow生成艺术化占位图（Mock模式）",
]
add_bullet_list(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1.5), gen_items, font_size=14, color=DARK, spacing=Pt(4))


# ========== 第7页：创意工坊 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RGBColor(0xE1, 0x70, 0x55))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
             "核心功能四：创意工坊", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), RGBColor(0xE1, 0x70, 0x55))
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(8), Inches(0.4),
             "多AI能力融合，实现跨模态创意组合", font_size=16, color=GRAY)

# 模式1：图片→文字
add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5), RGBColor(0xF0, 0xF0, 0xFF), PRIMARY)
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(0.5),
             "模式一：图片 → 文字", font_size=20, color=PRIMARY, bold=True)

mode1_items = [
    "Step 1: 用户上传图片",
    "Step 2: OpenCV 分析图像内容",
    "  - 提取颜色、检测人脸、识别场景",
    "Step 3: 将图像描述传递给 LLM",
    "Step 4: GLM-4 根据描述创作文字",
    "  - 可指定文学/技术等风格",
    "输出: 图像分析结果 + AI创作文字",
]
add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(3.5), mode1_items, font_size=13, color=DARK, spacing=Pt(4))

# 模式2：文字→图片
add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5), RGBColor(0xFF, 0xF0, 0xF5), ACCENT)
add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5), Inches(0.5),
             "模式二：文字 → 图片", font_size=20, color=ACCENT, bold=True)

mode2_items = [
    "Step 1: 用户输入创意描述",
    "Step 2: LLM 增强提示词",
    "  - 中文→英文，添加风格/光影/细节",
    "Step 3: 增强后的提示词传给 AIGC",
    "Step 4: CogView 生成精美图片",
    "  - 支持6种绘画风格",
    "输出: 增强提示词 + AI生成图片",
]
add_bullet_list(slide, Inches(7.4), Inches(2.8), Inches(4.8), Inches(3.5), mode2_items, font_size=13, color=DARK, spacing=Pt(4))


# ========== 第8页：数据库设计 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "数据库设计", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)

# 表格
tables_info = [
    ("UserProfile (用户资料)", [
        "id: Integer (PK)",
        "user: OneToOne(User)",
        "avatar: ImageField",
        "bio: TextField",
        "created_at: DateTime",
    ], PRIMARY),
    ("CreativeProject (创意项目)", [
        "id: Integer (PK)",
        "title: CharField(200)",
        "project_type: CharField(20)",
        "input_data: JSONField",
        "output_data: JSONField",
        "result_image: ImageField",
        "user: FK(User)",
        "is_public: BooleanField",
    ], SECONDARY),
    ("Artwork (作品)", [
        "id: Integer (PK)",
        "title: CharField(200)",
        "image: ImageField",
        "prompt: TextField",
        "ai_type: CharField(20)",
        "likes / views: Integer",
        "user: FK(User)",
    ], ACCENT),
    ("Comment / LikeRecord", [
        "Comment:",
        "  artwork: FK(Artwork)",
        "  user: FK(User)",
        "  content: TextField",
        "LikeRecord:",
        "  artwork+user: Unique",
    ], RGBColor(0xE1, 0x70, 0x55)),
]

table_left = Inches(0.8)
for title, fields, color in tables_info:
    card = add_shape(slide, table_left, Inches(1.6), Inches(2.9), Inches(5.0), WHITE, color)
    header = add_shape(slide, table_left, Inches(1.6), Inches(2.9), Inches(0.5), color)
    header_tf = header.text_frame
    header_tf.paragraphs[0].text = title
    header_tf.paragraphs[0].font.size = Pt(12)
    header_tf.paragraphs[0].font.color.rgb = WHITE
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    field_text = "\n".join(fields)
    add_text_box(slide, table_left + Inches(0.15), Inches(2.2), Inches(2.6), Inches(4.2),
                 field_text, font_size=10, color=DARK)
    table_left += Inches(3.1)


# ========== 第9页：项目结构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "项目结构", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)

structure = """期末项目/
├── zhiJing/              # Django项目配置
│   ├── settings.py       # 全局配置
│   ├── urls.py           # 根路由
│   └── wsgi.py
├── accounts/             # 用户管理App
│   ├── models.py         # UserProfile模型
│   ├── views.py          # 注册/登录/个人中心
│   └── forms.py          # 表单验证
├── creative/             # 创意工坊App
│   ├── models.py         # CreativeProject模型
│   └── views.py          # 识物/写作/绘画/工坊
├── gallery/              # 作品画廊App
│   ├── models.py         # Artwork/Comment/LikeRecord
│   └── views.py          # 画廊/点赞/评论
├── ai_service/           # Flask AI微服务
│   ├── app.py            # Flask路由(4个API)
│   ├── cv_service.py     # OpenCV视觉分析
│   ├── llm_service.py    # GLM-4大模型
│   └── aigc_service.py   # CogView图像生成
├── templates/            # 前端模板(12个页面)
├── static/css/           # 渐变科技风CSS
├── run_ai_service.py     # Flask启动脚本
├── start.bat             # 一键启动脚本
└── manage.py             # Django管理脚本"""

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.5),
             structure, font_size=11, color=DARK)

# 右侧：关键数据
add_shape(slide, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.5), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(7.7), Inches(1.6), Inches(4.5), Inches(0.5),
             "项目规模", font_size=20, color=PRIMARY, bold=True)

stats = [
    "Django Apps: 3个 (accounts / creative / gallery)",
    "Flask API端点: 4个 (health / cv / llm / aigc)",
    "数据模型: 5个 (UserProfile / CreativeProject / Artwork / Comment / LikeRecord)",
    "前端页面: 12个模板页面",
    "AI功能: 3大核心能力 + 1个融合模式",
    "写作风格: 5种",
    "绘画风格: 6种",
    "代码文件: 20+ 个Python文件",
    "技术栈: Django + Flask + OpenCV + ZhipuAI",
]
add_bullet_list(slide, Inches(7.9), Inches(2.2), Inches(4.3), Inches(4.5), stats, font_size=12, color=DARK, spacing=Pt(6))


# ========== 第10页：创新亮点 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "创新亮点", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

highlights = [
    ("01", "Flask + Django 微服务架构", "Django处理业务逻辑与页面渲染，Flask专注AI推理，职责分离、独立部署、易于扩展", PRIMARY),
    ("02", "多AI能力融合", "CV + LLM + AIGC 三种AI技术组合，创意工坊实现跨模态创作（图片→文字、文字→图片）", SECONDARY),
    ("03", "优雅降级设计", "无API Key时自动切换Mock演示模式，CV功能本地运行无需API，保证项目始终可演示", ACCENT),
    ("04", "完整创作闭环", "创作→保存→发布→点赞→评论，形成社区生态，用户可管理项目并发布到画廊", RGBColor(0xFD, 0xCB, 0x6E)),
]

card_top = Inches(1.6)
for num, title, desc, color in highlights:
    card = add_shape(slide, Inches(1.5), card_top, Inches(10.3), Inches(1.2), RGBColor(0x1A, 0x1A, 0x3E), color)

    num_shape = add_shape(slide, Inches(1.7), card_top + Inches(0.2), Inches(0.7), Inches(0.7), color)
    num_tf = num_shape.text_frame
    num_tf.paragraphs[0].text = num
    num_tf.paragraphs[0].font.size = Pt(20)
    num_tf.paragraphs[0].font.color.rgb = WHITE
    num_tf.paragraphs[0].font.bold = True
    num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.7), card_top + Inches(0.15), Inches(8.5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.7), card_top + Inches(0.6), Inches(8.5), Inches(0.5),
                 desc, font_size=13, color=RGBColor(0xA0, 0xA0, 0xC0))
    card_top += Inches(1.4)


# ========== 第11页：演示指南 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
top_bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "演示指南", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06), PRIMARY)

# 启动步骤
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "启动方式", font_size=20, color=PRIMARY, bold=True)

start_items = [
    "方式1：双击 start.bat 一键启动",
    "方式2：分别启动两个服务",
    "  终端1: python run_ai_service.py",
    "  终端2: python manage.py runserver 0.0.0.0:8000",
    "访问: http://localhost:8000",
]
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.2), Inches(2.5), start_items, font_size=14, color=DARK, spacing=Pt(6))

# 演示流程
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "推荐演示流程", font_size=20, color=ACCENT, bold=True)

demo_items = [
    "1. 首页浏览项目介绍",
    "2. 注册/登录账户",
    "3. 智能识物：上传图片查看分析",
    "4. AI写作：输入主题，切换风格",
    "5. AI绘画：描述画面，选择风格",
    "6. 创意工坊：体验多AI融合",
    "7. 作品画廊：浏览/点赞/评论",
]
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(3.0), demo_items, font_size=14, color=DARK, spacing=Pt(5))

# API配置
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0), LIGHT_BG, RGBColor(0xE0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(5.1), Inches(11), Inches(0.4),
             "API配置说明", font_size=18, color=PRIMARY, bold=True)

api_items = [
    "智谱AI API Key 已配置在 ai_service/config.py 中",
    "GLM-4-Flash 和 CogView-3-Flash 均为免费模型，新用户有免费额度",
    "OpenCV功能无需API，本地运行",
    "未配置API Key时，AI写作和绘画自动使用Mock演示模式",
]
add_bullet_list(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(1.3), api_items, font_size=13, color=DARK, spacing=Pt(4))


# ========== 第12页：总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1.0),
             "谢谢观看", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
             "智境 · AI创意工坊", font_size=28, color=RGBColor(0xA0, 0xA0, 0xD0), alignment=PP_ALIGN.CENTER)

summary_items = [
    "Flask + Django 微服务架构，职责分离",
    "OpenCV + GLM-4 + CogView 三大AI能力",
    "智能识物 / AI写作 / AI绘画 / 创意工坊",
    "完整用户系统与社区功能",
    "优雅降级，无API Key也可演示",
]
add_bullet_list(slide, Inches(3.0), Inches(3.3), Inches(7), Inches(2.5), summary_items, font_size=16, color=RGBColor(0xC0, 0xC0, 0xE0), spacing=Pt(10))

# 技术标签
tags = ["Flask", "Django", "OpenCV", "GLM-4", "CogView", "REST API"]
tag_left = Inches(3.0)
for tag in tags:
    tag_shape = add_shape(slide, tag_left, Inches(5.8), Inches(1.1), Inches(0.4), PRIMARY)
    tag_tf = tag_shape.text_frame
    tag_tf.paragraphs[0].text = tag
    tag_tf.paragraphs[0].font.size = Pt(12)
    tag_tf.paragraphs[0].font.color.rgb = WHITE
    tag_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_left += Inches(1.25)

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "Python Web 期末项目 | 2026年6月", font_size=14, color=RGBColor(0x80, 0x80, 0xA0), alignment=PP_ALIGN.CENTER)


# 保存
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "智境AI创意工坊_项目汇报.pptx")
prs.save(output_path)
print(f"PPT已生成: {output_path}")
